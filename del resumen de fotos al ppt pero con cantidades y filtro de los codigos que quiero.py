import os
import re
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ✅ Lista de códigos permitidos (copiados desde tu imagen)
codigos_filtrados = [
    "QD5154342",
    "QD5246801",
    "QD5154400",
    "F3786001",
    "F3786020",
    "QD5044460",
    "F3787001",
    "QD3953110",
    "F3786100",
    "F3787020",
    "F3787100",
    "QD5282001",
    "QD5045283",
    "QD5043332",
    "QD5043442"
]

# Rutas
input_path = r'C:\Users\Jaime Valderrama\OneDrive - American Sportswear, S.A\Documentos\Jaime\Testing\Nueva carpeta\Resultado\fotoresumen'
excel_path = r'C:\Users\Jaime Valderrama\OneDrive - American Sportswear, S.A\Documentos\Jaime\Testing\Nueva carpeta\fotos unides CH4\CK 1\cant.xlsx'
output_pptx = os.path.join(input_path, "ResumenPorReferencia_conCantidades_filtrado.pptx")

# Leer Excel sin encabezados
df = pd.read_excel(excel_path, header=None)
df.columns = ["Codigo", "Cantidad"]
df["Codigo"] = df["Codigo"].astype(str).str.strip()
cantidad_dict = dict(zip(df["Codigo"], df["Cantidad"]))

# Agrupar imágenes por referencia base
pattern = re.compile(r'^([A-Z0-9]+)_([0-9]+)_([0-9]+)\.(jpg|jpeg|png)$', re.IGNORECASE)

referencias = {}

for filename in os.listdir(input_path):
    match = pattern.match(filename)
    if match:
        ref, color, _, ext = match.groups()
        full_code = f"{ref}{color}"

        # Solo incluir si está en la lista filtrada (o todos si lista está vacía)
        if codigos_filtrados and full_code not in codigos_filtrados:
            continue

        img_path = os.path.join(input_path, filename)
        if ref not in referencias:
            referencias[ref] = []
        referencias[ref].append((img_path, full_code))

# Crear presentación
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

for ref, items in referencias.items():
    slide = prs.slides.add_slide(blank_slide_layout)

    # Título de la referencia
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = ref
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    x_start = Inches(0.5)
    y_start = Inches(1.2)
    img_width = Inches(2.3)
    spacing = Inches(0.3)
    max_per_row = 4
    row_height = Inches(3)

    for i, (img_path, full_code) in enumerate(items):
        col = i % max_per_row
        row = i // max_per_row

        if row >= 2:
            slide = prs.slides.add_slide(blank_slide_layout)
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = ref + " (cont.)"
            title_frame.paragraphs[0].font.size = Pt(28)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            y = y_start
            x = x_start
            row = 0
            col = 0

        x = x_start + col * (img_width + spacing)
        y = y_start + row * (row_height)

        slide.shapes.add_picture(img_path, x, y, width=img_width)

        cantidad = cantidad_dict.get(full_code, "Sin datos")
        subtitulo = f"{full_code} - {cantidad} unidades" if isinstance(cantidad, (int, float)) else f"{full_code} - {cantidad}"

        text_box = slide.shapes.add_textbox(x, y + Inches(2.4), img_width, Inches(0.4))
        tf = text_box.text_frame
        tf.text = subtitulo
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Guardar presentación
prs.save(output_pptx)
print(f"✅ Presentación generada SOLO con los códigos filtrados: {output_pptx}")
