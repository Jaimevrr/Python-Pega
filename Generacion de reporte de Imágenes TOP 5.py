import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR

# ======================================================
# CONFIGURACIÓN
# ======================================================

input_path = r'C:\Users\Jaime Valderrama\OneDrive - American Sportswear, S.A\Documentos\Jaime\Testing\Nueva carpeta\Resultado\Fotos 2025'
excel_path = r'C:\Users\Jaime Valderrama\OneDrive - American Sportswear, S.A\Documentos\Jaime\Testing\Nueva carpeta\1\ranking2025.xlsx'
template_path = r'C:\Users\Jaime Valderrama\OneDrive - American Sportswear, S.A\Documentos\Jaime\Testing\Nueva carpeta\1\Más Vendidos temp.pptx'
output_pptx = os.path.join(input_path, "Resumen2025EC.pptx")

TIPO_TIENDA_ACTIVO = "Internet"
ORDEN_GENERO = ["Men", "Women"]
USD_DIVISOR = 915

BASE_LAYOUT = 6  # SLIDE EN BLANCO (SIN TÍTULO)

# ======================================================
# LAYOUTS DE IMÁGENES
# ======================================================

LAYOUTS = {
    2: [(1.2, 1.8, 4.2, 3.0), (6.2, 1.8, 4.2, 3.0)],
    3: [(0.9, 1.8, 3.4, 2.9), (4.3, 1.8, 3.4, 2.9), (7.7, 1.8, 3.4, 2.9)],
    4: [(0.9, 1.7, 2.8, 2.7), (3.9, 1.7, 2.8, 2.7),
        (6.9, 1.7, 2.8, 2.7), (9.9, 1.7, 2.8, 2.7)],
    5: [(0.7, 1.7, 2.2, 2.5), (3.0, 1.7, 2.2, 2.5),
        (5.3, 1.7, 2.2, 2.5), (7.6, 1.7, 2.2, 2.5),
        (9.9, 1.7, 2.2, 2.5)]
}

TABLE_LEFT, TABLE_TOP = 1.2, 4.9
TABLE_WIDTH, TABLE_HEIGHT = 11.0, 1.3

# ======================================================
# FUNCIONES
# ======================================================

def buscar_imagen(ref):
    for ext in [".jpg", ".jpeg", ".png", ".webp"]:
        path = os.path.join(input_path, f"{ref}{ext}")
        if os.path.exists(path):
            return path
    return None


def slide_titulo(prs, texto):
    slide = prs.slides.add_slide(prs.slide_layouts[BASE_LAYOUT])

    # Caja centrada
    box = slide.shapes.add_textbox(
        Inches(2), Inches(2.7),
        Inches(9), Inches(1.4)
    )

    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # ✅ CORRECTO

    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER


def slide_productos(prs, df):
    n = min(max(len(df), 2), 5)
    slide = prs.slides.add_slide(prs.slide_layouts[BASE_LAYOUT])

    for i, (_, r) in enumerate(df.head(n).iterrows()):
        ref = str(r["Referencia"])
        img = buscar_imagen(ref)
        l, t, w, h = LAYOUTS[n][i]

        if img:
            slide.shapes.add_picture(img, Inches(l), Inches(t), Inches(w), Inches(h))

        sku = slide.shapes.add_textbox(
            Inches(l), Inches(t + h + 0.12), Inches(w), Inches(0.35)
        )
        p = sku.text_frame.paragraphs[0]
        p.text = ref
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

    tabla_df = df.head(n)[
        ["Referencia", "Descripción", "Suma de Cantidad", "Venta_USD"]
    ].copy()

    tabla_df["Venta_USD"] = tabla_df["Venta_USD"].apply(
        lambda x: f"USD {x:,.0f}".replace(",", ".")
    )

    table = slide.shapes.add_table(
        len(tabla_df) + 1, 4,
        Inches(TABLE_LEFT), Inches(TABLE_TOP),
        Inches(TABLE_WIDTH), Inches(TABLE_HEIGHT)
    ).table

    headers = ["Referencia", "Descripción", "Cantidad", "Venta (USD)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)

    for r, (_, row) in enumerate(tabla_df.iterrows(), start=1):
        table.cell(r, 0).text = row["Referencia"]
        table.cell(r, 1).text = row["Descripción"]
        table.cell(r, 2).text = str(int(row["Suma de Cantidad"]))
        table.cell(r, 3).text = row["Venta_USD"]

        for c in range(4):
            p = table.cell(r, c).text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.name = "Calibri"

# ======================================================
# MAIN
# ======================================================

def generar_ppt():
    df = pd.read_excel(excel_path)
    df.columns = df.columns.str.strip()

    df = df[df["Tipo de Tienda"] == TIPO_TIENDA_ACTIVO]
    df["Venta_USD"] = df["Suma de Total líneas"] / USD_DIVISOR

    df["Genero_Ord"] = df["Genero"].map({g: i for i, g in enumerate(ORDEN_GENERO)})
    df = df.sort_values(["Genero_Ord", "Categoria 2", "Ranking"])

    prs = Presentation(template_path)

    for genero in ORDEN_GENERO:
        for cat2, grp in df[df["Genero"] == genero].groupby("Categoria 2"):
            slide_titulo(prs, f"{cat2} — {genero}")
            slide_productos(prs, grp)

    prs.save(output_pptx)
    print("✅ PPT generado limpio, elegante y SIN placeholders:", output_pptx)


if __name__ == "__main__":
    generar_ppt()
