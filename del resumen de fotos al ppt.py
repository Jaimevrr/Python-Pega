import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Ruta de las imágenes finales
input_path = r'C:\Users\Jaime Valderrama\OneDrive - American Sportswear, S.A\Documentos\Jaime\Testing\Nueva carpeta\Resultado\fotoresumen'
output_pptx = os.path.join(input_path, "ResumenPorReferencia.pptx")

# Agrupar imágenes por referencia base
pattern = re.compile(r'^([A-Z0-9]+)_([0-9]+)_([0-9]+)\.(jpg|jpeg|png)$', re.IGNORECASE)

referencias = {}

for filename in os.listdir(input_path):
    match = pattern.match(filename)
    if match:
        ref, color, _, ext = match.groups()
        full_code = f"{ref}{color}"
        img_path = os.path.join(input_path, filename)
        if ref not in referencias:
            referencias[ref] = []
        referencias[ref].append((img_path, full_code))

# Crear presentación
prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

for ref, items in referencias.items():
    slide = prs.slides.add_slide(blank_slide_layout)

    # Título principal de la slide (la referencia)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = ref
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Distribuir imágenes horizontalmente (hasta 4 por slide, luego otra slide si hay más)
    x_start = Inches(0.5)
    y_start = Inches(1.2)
    img_width = Inches(2.3)
    spacing = Inches(0.3)
    max_per_row = 4
    row_height = Inches(3)

    for i, (img_path, full_code) in enumerate(items):
        col = i % max_per_row
        row = i // max_per_row

        x = x_start + col * (img_width + spacing)
        y = y_start + row * (row_height)

        # Si se pasa de 2 filas (8 imágenes), crear otra slide
        if row >= 2:
            slide = prs.slides.add_slide(blank_slide_layout)
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
            title_frame = title_box.text_frame
            title_frame.text = ref + " (cont.)"
            title_frame.paragraphs[0].font.size = Pt(28)
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            y = y_start
            x = x_start
            row = 0
            col = 0

        # Agregar imagen
        slide.shapes.add_picture(img_path, x, y, width=img_width)

        # Agregar texto debajo de la imagen
        text_box = slide.shapes.add_textbox(x, y + Inches(2.4), img_width, Inches(0.4))
        tf = text_box.text_frame
        tf.text = full_code
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# Guardar presentación
prs.save(output_pptx)
print(f"✅ Presentación generada: {output_pptx}")
